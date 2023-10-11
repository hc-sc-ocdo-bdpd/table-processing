import pandas as pd
import os
from difflib import SequenceMatcher
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from docx.api import Document
import datacompy
import itertools
import re
import subprocess
import logging
from glob import glob
from pathlib import Path
from pptx import Presentation

import numpy as np
from skimage.measure import label, regionprops

from fuzzywuzzy import fuzz
from fuzzywuzzy import process

import docx

import difflib

from os.path import splitext
from collections import Counter

import time

logger = logging.getLogger()
fhandler = logging.FileHandler(filename='nb_log.log', mode='a')
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
fhandler.setFormatter(formatter)
logger.addHandler(fhandler)
logger.setLevel(logging.WARNING)

# ---------------------------------------
# extract tables from word documents
# ---------------------------------------

def get_tbls_docx(path):
    for rootDirectory, subDirectories, fileList in os.walk(path):
        # for files in directory
        for fileName in fileList:
            # get file extension
            file_ext = fileName.rsplit('.', 1)[-1].lower()
            # if file is a docx
            if file_ext == "docx":
                # get full path
                fullPath = os.path.join(rootDirectory, fileName)
                # get file name
                doc = fileName.rsplit('.', 1)[0].lower()
                # extract document and tables
                document = Document(fullPath)
                tables = document.tables
                # create excel file
                writer = pd.ExcelWriter('../../tbls/'+doc+'_tbls.xlsx')
                # for each table in doc
                for count, table in enumerate(document.tables):
                    df = pd.DataFrame()
                    # write rows of table to df
                    for row in table.rows:
                        text = [cell.text for cell in row.cells]
                        df = pd.concat([df, pd.DataFrame([text])], ignore_index=True)
                    # write df to excel file
                    df.to_excel(writer, sheet_name=doc[:20]+"_tbl_"+str(count))
                writer.close()

# ---------------------------------------
# extract tables from powerpoints
# ---------------------------------------

def get_tbls_pptx(path):
    for rootDirectory, subDirectories, fileList in os.walk(path):
        # for files in directory
        for fileName in fileList:
            # get file extension
            file_ext = fileName.rsplit('.', 1)[-1].lower()
            # if file is a docx
            if file_ext == "pptx":
                # get full path
                fullPath = os.path.join(rootDirectory, fileName)
                # get file name
                doc = fileName.rsplit('.', 1)[0].lower()
                print(fullPath)
                # create excel file
                writer = pd.ExcelWriter('../../tbls/'+doc+'_tbls.xlsx')
                # extract document and tables
                prs = Presentation(fullPath)
                # text_runs will be populated with a list of strings,
                # one for each text run in presentation
                text_runs = []
                for count, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if not shape.has_table:
                            continue  
                        df = pd.DataFrame()  
                        tbl = shape.table
                        row_count = len(tbl.rows)
                        col_count = len(tbl.columns)
                        for r in range(0, row_count):
                            text = []
                            for c in range(0, col_count):
                                text_str = ""
                                cell = tbl.cell(r,c)
                                paragraphs = cell.text_frame.paragraphs 
                                for paragraph in paragraphs:
                                    for run in paragraph.runs:
                                        text_runs.append(run.text)
                                        text_str += run.text
                                text.append(text_str)
                            df = pd.concat([df, pd.DataFrame([text])], ignore_index=True)
                        
                        df.to_excel(writer, sheet_name=doc[:20]+"_tbl_"+str(count))
                writer.close()

# ---------------------------------------
# extract tables from xlsx
# ---------------------------------------

def get_tbls_xlsx(path):
    for rootDirectory, subDirectories, fileList in os.walk(path):
        # for files in directory
        for fileName in fileList:
            # get file extension
            file_ext = fileName.rsplit('.', 1)[-1].lower()
            # if file is a docx
            if file_ext == "xlsx":
                # get full path
                fullPath = os.path.join(rootDirectory, fileName)
                # get file name
                doc = fileName.rsplit('.', 1)[0].lower()
                # create excel file
                writer = pd.ExcelWriter('../../tbls/'+doc+'_tbls.xlsx')

                sheets_dict = pd.read_excel(fullPath, sheet_name=None, header=None)

                all_sheets = []
                count = 0

                for name, sheet in sheets_dict.items():
                    larr = label(np.array(sheet.notnull()).astype("int"))

                    for s in regionprops(larr):
                        try:
                            sub_df = (sheet.iloc[s.bbox[0]:s.bbox[2], s.bbox[1]:s.bbox[3]]
                                        .pipe(lambda sheet_: sheet_.rename(columns=sheet_.iloc[0])
                                        .drop(sheet_.index[0])))
                            
                            count +=1
                                        
                            sub_df.to_excel(writer, sheet_name=doc[:20]+"_tbl_"+str(count))
                        except Exception as e:
                            logging.info(e)
                writer.close()

# ---------------------------------------
# save tables from word documents into dictionary of dataframes
# ---------------------------------------

def convert_to_dict(base_path, doc):
    # get full path
    fullPath = os.path.join(base_path, doc)
    # create dictionary of dataframes
    doc_dict = pd.read_excel(fullPath, sheet_name=None)
    return doc_dict

# ---------------------------------------
# get paragraphs from docx to compare to tables
# ---------------------------------------

def getText(base_path, doc):
    fullPath = os.path.join(base_path, doc)
    docName = doc.rsplit('.', 1)[0].lower()
    document = docx.Document(fullPath)
    fullText = []
    for para in document.paragraphs:
        if len(para.text) > 50:
            fullText.append(para.text)
    text_name = [docName+"_para_"+str(x) for x in range(1,len(fullText))]
    # fullText2 = fullText
    # d1 = zip(fullText, fullText2)
    # fullTextDict = dict(d1)
    d1 = zip(text_name, fullText)
    fullTextDict = dict(d1)
    # fullTextDict = {v: k for v, k in enumerate(fullText)}
    return fullTextDict

# ---------------------------------------
# functions to get file names and match thresholds
# ---------------------------------------

def replace_para(short_name):
    short_name = str(short_name).split('_para')[0]
    matching = [s for s in fileNames if short_name.lower() in s.lower()[:len(short_name)]]
    try:
        return matching[0]
    except Exception as e:
        logging.info(e)
        return None

# set match thresholds
def tbl_match(row):
     # sum content and format match percentages
     if row["Compare_content"] + row["Compare_format"] > 1.25:
          return True
     else:
          return False

def replace(short_name):
    short_name = str(short_name).split('_tbl')[0]
    matching = [s for s in fileNames if short_name.lower() in s.lower()[:len(short_name)]]
    try:
        return matching[0]
    except Exception as e:
        logging.info(e)
        return None

# set match thresholds
def para_match(row):
     # sum content and format match percentages
     if row["Compare_content"] > 0.9:
          return True
     else:
          return False

# ---------------------------------------
# get list of comparisons for tables and paragraphs
# ---------------------------------------
def get_comparisons(stream):
    comparisons_df = pd.read_excel("comparisons.xlsx")

    tox_comp = comparisons_df[comparisons_df['Output Folder'] == stream]
    # tox_comp = tox_comp[tox_comp['File B'] == 3307206]
    tox_comp = tox_comp.drop('Output Folder', axis = 1)

    # get table comparisons
    # tox_comp['File A'] = tox_comp['File A'].astype(str) + '_tbls.xlsx'
    # tox_comp['File B'] = tox_comp['File B'].astype(str) + '_tbls.xlsx'


    # get paragraph comparisons
    tox_comp['File A'] = tox_comp['File A'].astype(str) + '.docx'
    tox_comp['File B'] = tox_comp['File B'].astype(str) + '.docx'
    tox_comp = tox_comp.drop_duplicates()
    comparisons = list(tox_comp.itertuples(index=False, name=None))
    return comparisons 

def compare_tbls():
    # path to tables from documents
    base_path = "C:\\Users\\EFUNG\\OneDrive - HC-SC PHAC-ASPC\\Documents\\my work\\PMRA doc disclosure\\tbls_F\\"
    timeout = 1000

    # for every pair of documents
    for pair in comparisons:
        timeout_start = time.time()
        # set up output dataframe
        dict_list = []
        df_name = "../../output_tbls_F/freas/"
        tbl_a = []
        tbl_b = []
        tbl_a_name = []
        tbl_b_name = []
        comp_val = []
        comp_val2 = []
        no_dict = False
        # for each document
        for doc in pair:
            df_name = df_name + str(doc)[:10] + "_"
            # save tables to dictionary
            try:
                doc_dict = convert_to_dict(base_path, doc)
                # add dictionary to list
                dict_list.append(doc_dict)
            except Exception as e:
                logging.info(e)
                no_dict = True
                print(doc)
                pass
        # compare every table in one dictionary to every table in the other dictionary
        if os.path.isfile(df_name + ".xlsx") == False:
            logging.info(df_name)
            if no_dict == False:
                logging.info("Tables")
                for a, b in itertools.combinations(dict_list, 2):
                    print(df_name)
                    for key_a in a:
                        for key_b in b:
                            try:
                                # while time.time() < timeout_start + timeout:
                                # compare tables - structure
                                compare = datacompy.Compare(a[key_a], b[key_b], abs_tol=0.5, rel_tol=0.5, on_index=True, ignore_case=True)
                                # get report on intersecting rows
                                report_df = compare.intersect_rows
                                # get columns with match results
                                match_cols = report_df.filter(regex='match')
                                # count the number of matches
                                true_count = (match_cols.values == True).sum()
                                # count the number of mismatches
                                false_count = (match_cols.values == False).sum()
                                # calculate percentages of matches
                                bool_comp = true_count / (true_count + false_count)                  
                                
                                # compare tables - content
                                # could try WRatio? Other ratios?
                                a_str = a[key_a].to_string()
                                b_str = b[key_b].to_string()
                                comp_val2.append(fuzz.WRatio(a_str, b_str)/100)
                                tbl_a_name.append(key_a)
                                tbl_b_name.append(key_b)
                                tbl_a.append(a_str)
                                tbl_b.append(b_str)
                                comp_val.append(bool_comp)
                            except Exception as e:
                                logging.info(e)

                # save match tables to excel file
                df_tbls = pd.DataFrame({'A_no': tbl_a_name, 'A': tbl_a, 'B_no': tbl_b_name, 'B': tbl_b, 'Compare_format': comp_val, 'Compare_content': comp_val2})
                # set match value
                df_tbls["Match"] = df_tbls.apply(lambda row: tbl_match(row), axis=1)
                # remove unnamed columns
                # cols = [c for c in df_tbls.columns if (c.lower()[:7] != 'unnamed') & (c.lower()[:6] != 'column')]
                # df_tbls=df_tbls[cols]

                df_tbls['file A'] = df_tbls.apply(lambda row: replace(row['A_no']), axis=1)
                df_tbls['file B'] = df_tbls.apply(lambda row: replace(row['B_no']), axis=1)

                df_name = df_name + ".xlsx"
                df_tbls.to_excel(df_name) 

def compare_para():
    doc_path = "C:\\Users\\EFUNG\\OneDrive - HC-SC PHAC-ASPC\\Documents\\my work\\PMRA doc disclosure\\FREAS\\"

    # for every pair of documents
    for pair in comparisons:
        # set up output dataframe
        dict_list = []
        df_name = "../../output_para_F/freas/"
        tbl_a = []
        tbl_b = []
        tbl_a_name = []
        tbl_b_name = []
        comp_val2 = []
        no_dict = False
        # for each document
        for doc in pair:
            df_name = df_name + str(doc)[:7] + "_"
            # save tables to dictionary
            try:
                doc_text = getText(doc_path, doc)
                # add dictionary to list
                dict_list.append(doc_text)
            except Exception as e:
                logging.info(e)
                no_dict = True
                print(doc)
                pass
        # compare every table in one dictionary to every table in the other dictionary
        if os.path.isfile(df_name + ".xlsx") == False:
            logging.info(df_name)
            if no_dict == False:
                logging.info("Paragraphs")
                for a, b in itertools.combinations(dict_list, 2):
                    for key_a in a:
                        for key_b in b:
                            try:
                                # compare paragraphs - content
                                # could try WRatio? Other ratios?
                                a_str = a[key_a]
                                b_str = b[key_b]
                                comp_val2.append(fuzz.WRatio(a_str, b_str)/100)
                                tbl_a_name.append(key_a)
                                tbl_b_name.append(key_b)
                                tbl_a.append(a_str)
                                tbl_b.append(b_str)
                            except Exception as e:
                                logging.info(e)

                # save match tables to excel file
                df_tbls = pd.DataFrame({'A_no': tbl_a_name, 'A': tbl_a, 'B_no': tbl_b_name, 'B': tbl_b, 'Compare_content': comp_val2})
                try:
                    # set match value
                    df_tbls["Match"] = df_tbls.apply(lambda row: para_match(row), axis=1)
                    # remove unnamed columns
                    # cols = [c for c in df_tbls.columns if (c.lower()[:7] != 'unnamed') & (c.lower()[:6] != 'column')]
                    # df_tbls=df_tbls[cols]

                    df_tbls['file A'] = df_tbls.apply(lambda row: replace(row['A_no']), axis=1)
                    df_tbls['file B'] = df_tbls.apply(lambda row: replace(row['B_no']), axis=1)

                    df_name = df_name + ".xlsx"
                    df_tbls.to_excel(df_name)
                except Exception as e:
                    print(df_name)
                    logging.info(e)

# ---------------------------------------
# set matches based on threshold
# ---------------------------------------

def get_matches():
    # path to output tables
    path = "C:\\Users\\EFUNG\\OneDrive - HC-SC PHAC-ASPC\\Documents\\my work\\PMRA doc disclosure\\output_tbls_F"

    df_all_matches = pd.DataFrame()

    for rootDirectory, subDirectories, fileList in os.walk(path):
        # for files in directory
        for fileName in fileList:
            fullPath = os.path.join(rootDirectory, fileName)
            df_tbls = pd.read_excel(fullPath)
            # df_tbls['file A'] = df_tbls.apply(lambda row: replace(row['A_no']), axis=1)
            # df_tbls['file B'] = df_tbls.apply(lambda row: replace(row['B_no']), axis=1)
            # df_tbls.to_excel(fullPath)
            # save matches to main match df
            to_append = df_tbls.loc[df_tbls["Match"] == True]
            # df_all_matches = df_all_matches.append(df_tbls.loc[df_tbls["Match"] == True])
            df_all_matches = pd.concat([df_all_matches, to_append])
            # save to excel
            # df_tbls.to_excel(fullPath)

    # save to excel
    df_all_matches.to_excel("df_all_matches_F.xlsx")